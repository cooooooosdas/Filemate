"""PPT 解析。"""

from __future__ import annotations

import logging
from pathlib import Path

from ..parsers import register

logger = logging.getLogger(__name__)


class PPTParser:
    """PPT 解析器。

    - .pptx → python-pptx（支持）
    - .ppt  → 旧格式，python-pptx 不支持
    """

    def parse(self, path: str | Path) -> dict:
        p = Path(path)
        suffix = p.suffix.lstrip(".").lower()

        if suffix == "ppt":
            return {
                "raw_text": "",
                "metadata": {"suffix": "ppt", "note": "旧 .ppt 格式暂不支持，请另存为 .pptx"},
                "error": "旧 .ppt 格式暂不支持，请另存为 .pptx 后重试",
            }

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx 未安装。运行: pip install python-pptx"
            ) from exc

        prs = Presentation(str(p))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- 第 {i} 页 ---\n" + "\n".join(texts))

        raw = "\n\n".join(slides)
        logger.debug("PPT 解析: %s → %d 页 / %d 字", p.name, len(slides), len(raw))
        return {
            "raw_text": raw,
            "metadata": {"suffix": "pptx", "slides": len(slides)},
        }


register("pptx", PPTParser)
register("PPTX", PPTParser)
